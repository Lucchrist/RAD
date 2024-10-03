import os
os.environ['HF_HUB_DISABLE_SYMLINKS_WARNING'] = '1'

import streamlit as st
import PyPDF2
import pptx
from docx import Document
from transformers import pipeline

# Charger le modèle BART pour résumer les documents
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Fonction pour lire jusqu'à 9 pages d'un fichier PDF
def read_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    max_pages = min(9, len(reader.pages))  # Limite à 9 pages ou au nombre de pages réel s'il y en a moins
    for page_num in range(max_pages):
        text += reader.pages[page_num].extract_text()
    return text

# Fonction pour tronquer le texte si nécessaire
def truncate_content(content, max_length=1024):
    return content[:max_length] if len(content) > max_length else content

# Fonction pour lire un fichier Word
def read_word(file):
    doc = Document(file)
    text = [paragraph.text for paragraph in doc.paragraphs]
    return '\n'.join(text)

# Fonction pour lire un fichier PowerPoint
def read_ppt(file):
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

# Interface Streamlit
st.title("Téléverser et résumer un document PDF, Word ou PowerPoint (Limité à 9 pages pour PDF)")

# Téléversement du fichier
uploaded_file = st.file_uploader("Choisissez un fichier (PDF, Word, PowerPoint)", type=["pdf", "docx", "pptx"])

# Afficher le contenu du fichier téléversé
if uploaded_file is not None:
    file_type = uploaded_file.name.split('.')[-1]
    
    # Lire et afficher le fichier en fonction de son type
    if file_type == "pdf":
        st.subheader("Visualisation du contenu PDF (limité à 9 pages)")
        content = read_pdf(uploaded_file)
        st.text(content[:2000])  # Limite l'affichage à 2000 caractères
    
    elif file_type == "docx":
        st.subheader("Visualisation du contenu Word")
        content = read_word(uploaded_file)
        st.text(content[:2000])  # Limite l'affichage à 2000 caractères
    
    elif file_type == "pptx":
        st.subheader("Visualisation du contenu PowerPoint")
        content = read_ppt(uploaded_file)
        st.text(content[:2000])  # Limite l'affichage à 2000 caractères

    # Options pour choisir le type de résumé
    st.subheader("Obtenir un résumé")
    
    # Tronquer le texte avant de le résumer
    truncated_content = truncate_content(content)
    
    text_length = len(truncated_content.split())
    summary_length_half = text_length // 2
    summary_length_third = text_length // 3
    
    if st.button("Résumé abrégé (1/3 de la longueur)"):
        summary = summarizer(truncated_content, max_length=summary_length_third, min_length=summary_length_third - 50, do_sample=False, clean_up_tokenization_spaces=True)
        st.subheader("Résumé abrégé")
        st.write(summary[0]['summary_text'])
    
    if st.button("Résumé détaillé (1/2 de la longueur)"):
        summary = summarizer(truncated_content, max_length=summary_length_half, min_length=summary_length_half - 50, do_sample=False, clean_up_tokenization_spaces=True)
        st.subheader("Résumé détaillé")
        st.write(summary[0]['summary_text'])
